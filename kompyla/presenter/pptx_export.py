"""Render a markdown document to .pptx — one slide per `## ` heading."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


_HEADING_RE = re.compile(r"^(#{1,3})\s+(.+)$")
_BULLET_RE = re.compile(r"^\s*[-*]\s+(.*)$")


def _strip_inline(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"\[\[([^\]]+)\]\]", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text


def _split_slides(markdown_text: str) -> list[tuple[str, list[str]]]:
    """Return [(slide_title, [bullet_lines, ...]), ...]."""
    lines = markdown_text.splitlines()
    slides: list[tuple[str, list[str]]] = []
    cover_title: str | None = None
    cover_body: list[str] = []
    i = 0
    in_fence = False

    # Skip frontmatter
    if lines and lines[0].startswith("---"):
        i = 1
        while i < len(lines) and not lines[i].startswith("---"):
            i += 1
        i += 1

    # Cover slide content (until first H2)
    while i < len(lines):
        line = lines[i]
        if line.startswith("## "):
            break
        if line.startswith("```"):
            in_fence = not in_fence
            i += 1
            continue
        if in_fence:
            i += 1
            continue
        if m := _HEADING_RE.match(line):
            if m.group(1) == "#" and not cover_title:
                cover_title = _strip_inline(m.group(2))
        elif line.strip():
            cover_body.append(_strip_inline(line.strip()))
        i += 1
    slides.append((cover_title or "Title", cover_body))

    # Body slides
    current_title: str | None = None
    current_body: list[str] = []
    while i < len(lines):
        line = lines[i]
        if line.startswith("```"):
            in_fence = not in_fence
            i += 1
            continue
        if not in_fence and line.startswith("## "):
            if current_title is not None:
                slides.append((current_title, current_body))
            current_title = _strip_inline(line[3:].strip())
            current_body = []
        elif current_title is not None and line.strip():
            if m := _BULLET_RE.match(line):
                current_body.append(_strip_inline(m.group(1)))
            elif _HEADING_RE.match(line):
                continue   # skip sub-headings; keep slides flat
            else:
                current_body.append(_strip_inline(line.strip()))
        i += 1
    if current_title is not None:
        slides.append((current_title, current_body))

    return slides


def md_to_pptx(markdown_text: str, out_path: Path, *, deck_title: str | None = None) -> Path:
    prs = Presentation()
    blank = prs.slide_layouts[5]  # Title-only layout

    slides = _split_slides(markdown_text)
    if deck_title:
        slides.insert(0, (deck_title, []))

    for title, body in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = title

        if body:
            tx_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.6), Inches(8.5), Inches(5.0),
            )
            tf = tx_box.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(body[:12]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.level = 0

    prs.save(out_path)
    return out_path
